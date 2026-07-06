#!/usr/bin/env python3
"""
Generateur PowerPoint — Soutenance (RE)Sources Relationnelles
12 slides conformes au plan valide dans Notes_oral_et_slides.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Soutenance_RE_Sources.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Palette ─────────────────────────────────────────────────────────────────
PRIMARY      = RGBColor(0x1D, 0x3A, 0x6D)
PRIMARY_DARK = RGBColor(0x0D, 0x2A, 0x5D)
ACCENT       = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_GOLD  = RGBColor(0xFF, 0xC1, 0x07)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG     = RGBColor(0xF2, 0xF6, 0xFA)
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)
GRAY         = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xBB, 0xBB, 0xBB)
RED_SOFT     = RGBColor(0xE5, 0x3E, 0x3E)

FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"

# ─── Helper functions ────────────────────────────────────────────────────────

def solid_bg(slide, color):
    """Fill the slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def gradient_bg(slide, c1, c2):
    """Approximate gradient with a full-slide rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.gradient()
    shape.fill.gradient_stops[0].color.rgb = c1
    shape.fill.gradient_stops[0].position = 0.0
    shape.fill.gradient_stops[1].color.rgb = c2
    shape.fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()
    # push behind everything
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)

def add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                  font_name=FONT_BODY, italic=False, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(4)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     color=DARK_TEXT, bold_prefix=False, spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(6)
        p.line_spacing = Pt(font_size * spacing)
        p.level = 0
    return tf

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body, icon="",
             title_color=PRIMARY, bg_color=WHITE):
    """Add a rounded card with icon, title and body text."""
    card = add_shape_rect(slide, left, top, width, height, bg_color, border_color=RGBColor(0xDD,0xDD,0xDD), radius=True)
    # Icon
    if icon:
        add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.6), Inches(0.5),
                     icon, font_size=22, alignment=PP_ALIGN.CENTER)
    # Title
    x_offset = Inches(0.8) if icon else Inches(0.3)
    add_text_box(slide, left + x_offset, top + Inches(0.15), width - x_offset - Inches(0.2), Inches(0.4),
                 title, font_size=15, bold=True, color=title_color)
    # Body
    add_text_box(slide, left + Inches(0.3), top + Inches(0.55), width - Inches(0.6), height - Inches(0.7),
                 body, font_size=12, color=GRAY, line_spacing=1.3)

def slide_number_footer(slide, num, total=14):
    """Add a subtle slide number bottom-right."""
    add_text_box(slide, prs.slide_width - Inches(1.2), prs.slide_height - Inches(0.45),
                 Inches(1), Inches(0.35),
                 f"{num} / {total}", font_size=10, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.RIGHT)

def section_header(slide, number, title, subtitle="", timing=""):
    """Add a consistent section header bar at the top."""
    # Blue bar
    add_shape_rect(slide, 0, 0, prs.slide_width, Inches(1.15), PRIMARY)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.2), Inches(0.65), Inches(0.65))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    add_text_box(slide, Inches(1.4), Inches(0.15), Inches(9), Inches(0.5),
                 title, font_size=28, bold=True, color=WHITE, font_name=FONT_TITLE)
    # Subtitle
    if subtitle:
        add_text_box(slide, Inches(1.4), Inches(0.6), Inches(9), Inches(0.4),
                     subtitle, font_size=14, color=RGBColor(0xBB, 0xCC, 0xEE), italic=True)
    # Timing badge
    if timing:
        add_text_box(slide, prs.slide_width - Inches(2.2), Inches(0.35), Inches(1.8), Inches(0.4),
                     timing, font_size=11, color=RGBColor(0xAA, 0xBB, 0xDD),
                     alignment=PP_ALIGN.RIGHT, italic=True)

def placeholder_image(slide, left, top, width, height, label):
    """Insert a gray dashed placeholder for an image to be added later."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 4  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"[IMAGE A INSERER]\n{label}"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER


# =============================================================================
# SLIDE 1 — Page de titre
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
gradient_bg(slide, PRIMARY_DARK, PRIMARY)

# Central content
add_text_box(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(1),
             "(RE)Sources Relationnelles", font_size=48, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

add_text_box(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(0.6),
             "Plateforme de ressources pour l'amelioration des relations entre citoyens",
             font_size=20, color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER, italic=True)

# Separator line
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.2), Inches(2.3), Pt(3))
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_GOLD
sep.line.fill.background()

add_text_box(slide, Inches(1), Inches(3.6), Inches(11.3), Inches(0.5),
             "Projet collaboratif — INFCDAAL2 — CESI Ecole d'Ingenieurs",
             font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

# Team cards
team = [
    ("Chamil MUSAYEV", "Backend Lead", "API, BDD, Auth, Securite"),
    ("Jamal", "Frontend Lead", "Next.js, RGAA, Integration"),
    ("Romain CONDOMITTI", "Back-office", "Admin, Auth, Moderation"),
    ("Leon HEU", "Mobile Lead", "React Native, Expo"),
]
card_w = Inches(2.7)
card_h = Inches(1.35)
start_x = Inches(0.85)
y = Inches(4.6)
gap = Inches(0.35)

for i, (name, role, scope) in enumerate(team):
    x = start_x + i * (card_w + gap)
    card = add_shape_rect(slide, x, y, card_w, card_h, RGBColor(0x15, 0x2E, 0x55), radius=True)
    card.line.color.rgb = RGBColor(0x3A, 0x5A, 0x8A)
    card.line.width = Pt(1)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.12), card_w - Inches(0.3), Inches(0.35),
                 name, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.48), card_w - Inches(0.3), Inches(0.3),
                 role, font_size=12, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.82), card_w - Inches(0.3), Inches(0.35),
                 scope, font_size=10, color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

# Bottom date
add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
             "Soutenance — Avril 2026", font_size=14, color=RGBColor(0x88, 0x99, 0xBB),
             alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 1)


# =============================================================================
# SLIDE 2 — Contexte et enjeux
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 2, "Contexte et enjeux", "Le Ministere des Solidarites et de la Sante", "1 min 30")

# Left : text content
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(6.5), Inches(0.5),
             "Pourquoi ce projet ?", font_size=22, bold=True, color=PRIMARY)

bullets = [
    "Commande simulee du Ministere des Solidarites et de la Sante",
    "Constat : la qualite de nos relations est un levier majeur de bien-etre (Maslow)",
    "Objectif : plateforme de ressources pour tous les types de relations",
    "Cible : couple, famille, amis, collegues — tous les citoyens",
    "Fonctionnalites : consultation, creation, echanges moderes, progression",
]
add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(6.3), Inches(3.5),
                [f"   {b}" for b in bullets], font_size=16, color=DARK_TEXT)

# Right : image placeholder for Maslow pyramid
placeholder_image(slide, Inches(7.8), Inches(1.5), Inches(4.8), Inches(3.5),
                  "Pyramide de Maslow\n(schema simplifie)")

# Bottom key message
msg_box = add_shape_rect(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.9), RGBColor(0xEA, 0xF4, 0xFB), radius=True)
add_text_box(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.7),
             "\"L'un des leviers les plus importants reste la qualite de nos relations aux autres\" — Cahier des charges CCTP",
             font_size=15, italic=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 2)


# =============================================================================
# SLIDE 3 — Analyse fonctionnelle
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 3, "Analyse fonctionnelle", "Diagramme de cas d'utilisation", "1 min 30")

# Full-width image placeholder for Use Case diagram
placeholder_image(slide, Inches(0.6), Inches(1.4), Inches(7.5), Inches(5.3),
                  "Diagramme de cas d'utilisation UML\n\n5 acteurs : Citoyen non connecte, Citoyen connecte,\nModerateur, Administrateur, Super-admin\n\n(draw.io / PlantUML / StarUML)")

# Right side: actor summary
actors = [
    ("Citoyen", "Consulte, cree, commente, favoris"),
    ("Moderateur", "Valide ressources et commentaires"),
    ("Administrateur", "Statistiques, suspension, catalogue"),
    ("Super-admin", "Cree comptes privilegies"),
]
y_start = Inches(1.6)
for i, (actor, desc) in enumerate(actors):
    y = y_start + i * Inches(1.15)
    add_card(slide, Inches(8.5), y, Inches(4.2), Inches(0.95),
             actor, desc, icon=["👤", "🛡️", "📊", "👑"][i])

slide_number_footer(slide, 3)


# =============================================================================
# SLIDE 4 — Analyse de la valeur
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 4, "Analyse de la valeur", "Bete a cornes et diagramme pieuvre", "1 min")

# Two image placeholders side by side
placeholder_image(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(4.2),
                  "Diagramme bete a cornes\n\nA qui ? Citoyens, familles, couples\nSur quoi ? Relations interpersonnelles\nPourquoi ? Ameliorer la qualite des liens")

placeholder_image(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.2),
                  "Diagramme pieuvre\n\nFP1 : Acces web aux ressources\nFP2 : Acces mobile\nFC1 : RGPD / RGAA\nFC2 : Performance\nFC3 : Navigateurs")

# Bottom legend
add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
             "FP = Fonction Principale    |    FC = Fonction Contrainte    |    Les deux diagrammes ont ete realises conformement au CCTP",
             font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER, italic=True)

slide_number_footer(slide, 4)


# =============================================================================
# SLIDE 5 — Choix techniques
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 5, "Choix technologiques", "Pourquoi Laravel + Next.js + PostgreSQL ?", "1 min 30")

# Comparison table area
# Headers
cols_x = [Inches(0.4), Inches(3.6), Inches(6.7), Inches(9.8)]
cols_w = [Inches(2.8), Inches(2.8), Inches(2.8), Inches(2.8)]
header_y = Inches(1.5)
header_h = Inches(0.55)

headers = ["Critere", "Laravel 11 (retenu)", "Symfony 7", "Express.js"]
header_colors = [PRIMARY, ACCENT, RGBColor(0x78, 0x78, 0x78), RGBColor(0x78, 0x78, 0x78)]

for i, (hdr, col) in enumerate(zip(headers, header_colors)):
    rect = add_shape_rect(slide, cols_x[i], header_y, cols_w[i], header_h, col if i <= 1 else RGBColor(0xEE, 0xEE, 0xEE))
    add_text_box(slide, cols_x[i], header_y + Inches(0.08), cols_w[i], header_h,
                 hdr, font_size=13, bold=True,
                 color=WHITE if i <= 1 else DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Table rows
rows_data = [
    ("MVC natif", "Oui, convention over config", "Oui, plus verbeux", "Non, structuration manuelle"),
    ("Auth integree", "Sanctum (tokens API legers)", "Security Bundle (complexe)", "Passport.js (tiers)"),
    ("ORM", "Eloquent (Active Record)", "Doctrine (Data Mapper)", "Sequelize/Prisma (tiers)"),
    ("Validation", "FormRequest declares", "Annotations/constraints", "Joi/Yup (tiers)"),
    ("Courbe apprentissage", "Moderee, doc excellente", "Elevee", "Faible mais pas de structure"),
]

row_y = header_y + header_h + Inches(0.05)
row_h = Inches(0.6)

for r_idx, (crit, *vals) in enumerate(rows_data):
    bg = LIGHT_BG if r_idx % 2 == 0 else WHITE
    for c_idx, val in enumerate([crit] + list(vals)):
        rx = cols_x[c_idx]
        add_shape_rect(slide, rx, row_y, cols_w[c_idx], row_h, bg)
        fc = PRIMARY if c_idx == 0 else (ACCENT if c_idx == 1 else GRAY)
        fb = c_idx == 0
        add_text_box(slide, rx + Inches(0.1), row_y + Inches(0.05), cols_w[c_idx] - Inches(0.2), row_h - Inches(0.1),
                     val, font_size=11, color=fc, bold=fb)
    row_y += row_h

# Bottom summary cards
card_data = [
    ("Next.js 16", "App Router, SSR, TypeScript, Tailwind v4\nExpertise React de l'equipe"),
    ("PostgreSQL 16", "Chiffrement natif pour RGPD\nRobustesse ACID, image Docker legere"),
    ("React Native / Expo", "Cross-platform iOS + Android\nMeme ecosysteme que le web"),
]
card_y = Inches(5.4)
card_w_bottom = Inches(3.8)
for i, (title, body) in enumerate(card_data):
    x = Inches(0.6) + i * (card_w_bottom + Inches(0.45))
    add_card(slide, x, card_y, card_w_bottom, Inches(1.3), title, body,
             icon=["⚛️", "🐘", "📱"][i])

slide_number_footer(slide, 5)


# =============================================================================
# SLIDE 6 — Comparaison et choix des technologies (multi-couches)
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 6, "Choix technologiques — Analyse comparative",
               "Frontend / Backend / Base de donnees", "1 min 30")

# Sub-header explanation
add_text_box(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.35),
             "Evaluation multi-criteres sur les 3 couches principales de l'application",
             font_size=12, italic=True, color=GRAY)

# 3 columns: Frontend, Backend, Database
stacks = [
    {
        "layer": "Frontend Web",
        "color": RGBColor(0x61, 0xDA, 0xFB),  # React blue
        "choice": "Next.js 16 + React 19",
        "alt1":   "Vue 3 / Nuxt 3",
        "alt2":   "Angular 18",
        "scores": [
            ("Performance",          "★★★★★", "★★★★☆", "★★★★☆"),
            ("Prise en main",        "★★★★★", "★★★★★", "★★★☆☆"),
            ("Ecosysteme React",     "★★★★★", "★★☆☆☆", "★☆☆☆☆"),
            ("SSR / App Router",     "★★★★★", "★★★★☆", "★★★☆☆"),
            ("Licence",              "MIT",    "MIT",    "MIT"),
        ],
        "why": "Ecosysteme React maitrise, App Router + SSR natifs, TypeScript et Tailwind v4 de premier ordre.",
    },
    {
        "layer": "Backend API",
        "color": ACCENT,
        "choice": "Laravel 12 (PHP 8.2+)",
        "alt1":   "Symfony 7",
        "alt2":   "Express.js",
        "scores": [
            ("Performance",          "★★★★☆", "★★★★☆", "★★★★★"),
            ("Prise en main",        "★★★★★", "★★★☆☆", "★★★★☆"),
            ("MVC natif",            "★★★★★", "★★★★★", "★☆☆☆☆"),
            ("Auth integree",        "★★★★★", "★★★★☆", "★★☆☆☆"),
            ("Licence",              "MIT",    "MIT",    "MIT"),
        ],
        "why": "MVC convention-over-config, Eloquent ORM, Sanctum pour tokens API, documentation exceptionnelle.",
    },
    {
        "layer": "Base de donnees",
        "color": RGBColor(0x33, 0x67, 0x91),  # PostgreSQL blue
        "choice": "PostgreSQL 16",
        "alt1":   "MySQL 8",
        "alt2":   "MongoDB",
        "scores": [
            ("Performance",          "★★★★★", "★★★★☆", "★★★★☆"),
            ("Conformite SQL",       "★★★★★", "★★★★☆", "★☆☆☆☆"),
            ("Integrite ACID",       "★★★★★", "★★★★☆", "★★★☆☆"),
            ("Chiffrement natif",    "★★★★★", "★★★☆☆", "★★★☆☆"),
            ("Licence",              "PostgreSQL", "GPL dual", "SSPL"),
        ],
        "why": "Chiffrement natif adapte au RGPD, robustesse ACID, image Docker officielle legere.",
    },
]

col_w = Inches(4.15)
col_gap = Inches(0.15)
col_y = Inches(1.7)
col_h = Inches(4.4)

for i, s in enumerate(stacks):
    cx = Inches(0.45) + i * (col_w + col_gap)

    # Layer header
    hdr = add_shape_rect(slide, cx, col_y, col_w, Inches(0.5), s["color"], radius=True)
    add_text_box(slide, cx, col_y + Inches(0.07), col_w, Inches(0.4),
                 s["layer"], font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Technology choice box (highlighted)
    choice_y = col_y + Inches(0.6)
    choice_box = add_shape_rect(slide, cx, choice_y, col_w, Inches(0.55),
                                LIGHT_BG, border_color=s["color"], radius=True)
    choice_box.line.width = Pt(2)
    add_text_box(slide, cx + Inches(0.1), choice_y + Inches(0.04), col_w - Inches(0.2), Inches(0.25),
                 "RETENU", font_size=8, bold=True, color=s["color"],
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, cx + Inches(0.1), choice_y + Inches(0.22), col_w - Inches(0.2), Inches(0.3),
                 s["choice"], font_size=12, bold=True, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)

    # Mini table: criterion | retained | alt1 | alt2
    table_y = choice_y + Inches(0.7)
    row_h = Inches(0.32)
    sub_cols_w = [Inches(1.55), Inches(0.95), Inches(0.85), Inches(0.8)]
    sub_cols_x = [cx]
    for w in sub_cols_w[:-1]:
        sub_cols_x.append(sub_cols_x[-1] + w)

    # Sub-header
    sub_headers = ["Critere", s["choice"].split()[0], s["alt1"].split()[0], s["alt2"].split()[0]]
    for c_idx, hdr_text in enumerate(sub_headers):
        bg = s["color"] if c_idx == 1 else RGBColor(0xDD, 0xDD, 0xDD)
        fc = WHITE if c_idx == 1 else DARK_TEXT
        add_shape_rect(slide, sub_cols_x[c_idx], table_y, sub_cols_w[c_idx], row_h, bg)
        add_text_box(slide, sub_cols_x[c_idx] + Inches(0.02), table_y + Inches(0.04),
                     sub_cols_w[c_idx] - Inches(0.04), row_h,
                     hdr_text, font_size=9, bold=True, color=fc, alignment=PP_ALIGN.CENTER)

    # Scores rows
    for r_idx, row in enumerate(s["scores"]):
        ry = table_y + row_h + r_idx * row_h
        bg = LIGHT_BG if r_idx % 2 == 0 else WHITE
        for c_idx, val in enumerate(row):
            add_shape_rect(slide, sub_cols_x[c_idx], ry, sub_cols_w[c_idx], row_h, bg)
            fc = s["color"] if c_idx == 1 else DARK_TEXT
            fb = c_idx == 0 or c_idx == 1
            add_text_box(slide, sub_cols_x[c_idx] + Inches(0.03), ry + Inches(0.04),
                         sub_cols_w[c_idx] - Inches(0.06), row_h,
                         val, font_size=9, bold=fb, color=fc, alignment=PP_ALIGN.CENTER)

    # "Why this choice" note
    why_y = table_y + row_h * (len(s["scores"]) + 1) + Inches(0.1)
    add_text_box(slide, cx + Inches(0.05), why_y, col_w - Inches(0.1), Inches(0.6),
                 "Pourquoi ce choix ?", font_size=10, bold=True, color=s["color"])
    add_text_box(slide, cx + Inches(0.05), why_y + Inches(0.28), col_w - Inches(0.1), Inches(0.8),
                 s["why"], font_size=9, color=GRAY, italic=True, line_spacing=1.25)

# Bottom note
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
             "Comparaison etablie sur les criteres prioritaires du projet : RGPD, maintenabilite, competences equipe",
             font_size=10, italic=True, color=GRAY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 6)


# =============================================================================
# SLIDE 7 — Respect du pattern MVC
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 7, "Respect du pattern MVC",
               "Separation stricte des responsabilites", "1 min")

# Sub-header
add_text_box(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.35),
             "Architecture en couches decouplees : API Laravel (Model + Controller) + SPA Next.js (View)",
             font_size=12, italic=True, color=GRAY)

# Three MVC blocks
MVC_BLUE   = RGBColor(0x21, 0x96, 0xF3)   # Model = blue
MVC_PURPLE = RGBColor(0x7E, 0x57, 0xC2)   # Controller = purple
MVC_GREEN  = RGBColor(0x4C, 0xAF, 0x50)   # View = green

mvc_blocks = [
    {
        "role": "MODEL",
        "color": MVC_BLUE,
        "subtitle": "Donnees + logique metier",
        "files": [
            "app/Models/User.php",
            "app/Models/Resource.php",
            "app/Models/Comment.php",
            "app/Models/Favorite.php",
            "+ Category, Progression,\n  RelationType, ResourceType",
        ],
        "note": "Eloquent ORM\nRelations, scopes, casts",
    },
    {
        "role": "CONTROLLER",
        "color": MVC_PURPLE,
        "subtitle": "Requetes + routing + validation",
        "files": [
            "app/Http/Controllers/\n  AuthController.php",
            "ResourceController.php",
            "CommentController.php",
            "AdminController.php",
            "+ Moderation, Favorite,\n  Progression, SuperAdmin",
        ],
        "note": "FormRequest + Policies\nMiddleware auth:sanctum",
    },
    {
        "role": "VIEW",
        "color": MVC_GREEN,
        "subtitle": "Affichage + UI + rendu",
        "files": [
            "apps/web/src/app/\n  (main)/ressources/",
            "(main)/auth/connexion/",
            "administration/",
            "components/ui/Button, Card,\n  Badge, Input",
            "+ ResourceCard, Navbar,\n  AuthContext",
        ],
        "note": "Next.js 16 App Router\nReact 19 + TypeScript + Tailwind",
    },
]

block_w = Inches(3.85)
block_gap = Inches(0.35)
block_y = Inches(1.75)
block_h = Inches(3.7)

# Arrows between blocks
arrow_w = Inches(0.35)
arrow_h = Inches(0.45)

for i, b in enumerate(mvc_blocks):
    bx = Inches(0.5) + i * (block_w + block_gap)

    # Block background
    card = add_shape_rect(slide, bx, block_y, block_w, block_h,
                          WHITE, border_color=b["color"], radius=True)
    card.line.width = Pt(2.5)

    # Top band with role
    band = add_shape_rect(slide, bx, block_y, block_w, Inches(0.75), b["color"], radius=True)
    # Force band to be visually on top of the card top
    add_text_box(slide, bx, block_y + Inches(0.08), block_w, Inches(0.35),
                 b["role"], font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, bx, block_y + Inches(0.42), block_w, Inches(0.3),
                 b["subtitle"], font_size=10, italic=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # File list
    files_y = block_y + Inches(0.95)
    add_text_box(slide, bx + Inches(0.2), files_y, block_w - Inches(0.4), Inches(0.3),
                 "Fichiers reels du projet", font_size=9, bold=True, color=b["color"])

    for f_idx, f in enumerate(b["files"]):
        fy = files_y + Inches(0.32) + f_idx * Inches(0.38)
        # Small bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     bx + Inches(0.25), fy + Inches(0.08),
                                     Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = b["color"]
        dot.line.fill.background()
        add_text_box(slide, bx + Inches(0.4), fy, block_w - Inches(0.5), Inches(0.35),
                     f, font_size=10, color=DARK_TEXT, line_spacing=1.15)

    # Note at bottom of block
    note_y = block_y + block_h - Inches(0.75)
    add_shape_rect(slide, bx + Inches(0.15), note_y, block_w - Inches(0.3), Inches(0.6),
                   LIGHT_BG, radius=True)
    add_text_box(slide, bx + Inches(0.2), note_y + Inches(0.08), block_w - Inches(0.4), Inches(0.5),
                 b["note"], font_size=9, color=GRAY, italic=True,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.2)

# Arrows between blocks (bidirectional)
for i in range(2):
    ax = Inches(0.5) + (i + 1) * block_w + i * block_gap + Inches(-0.03)
    ay = block_y + block_h / 2 - arrow_h / 2
    arrow = slide.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, ax, ay, arrow_gap_w := arrow_w + Inches(0.06), arrow_h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
    arrow.line.fill.background()

# Conformity table below
table_y = Inches(5.7)
table_title_y = table_y - Inches(0.25)
add_text_box(slide, Inches(0.5), table_title_y, Inches(12.3), Inches(0.3),
             "Conformite MVC", font_size=13, bold=True, color=PRIMARY)

conf_rows = [
    ("Role MVC",   "Responsabilite",                             "Implementation",                                 "Conforme"),
    ("Model",      "Donnees, logique metier, acces BDD",         "8 modeles Eloquent (User, Resource, Comment...)", "Oui"),
    ("Controller", "Traitement HTTP, routing, validation",       "9 controllers + FormRequest + Policies",         "Oui"),
    ("View",       "Affichage, templates, rendu",                "SPA Next.js decouplee (App Router + composants)", "Oui"),
]

conf_x = [Inches(0.5), Inches(2.0), Inches(4.5), Inches(11.0)]
conf_w = [Inches(1.5), Inches(2.5), Inches(6.5), Inches(1.8)]
conf_rh = Inches(0.4)

for r_idx, row in enumerate(conf_rows):
    ry = table_y + r_idx * conf_rh
    if r_idx == 0:
        bg = PRIMARY
        fc = WHITE
        fb = True
    else:
        bg = LIGHT_BG if r_idx % 2 == 1 else WHITE
        fc = DARK_TEXT
        fb = False
    for c_idx, val in enumerate(row):
        add_shape_rect(slide, conf_x[c_idx], ry, conf_w[c_idx], conf_rh, bg,
                       border_color=RGBColor(0xDD, 0xDD, 0xDD))
        align = PP_ALIGN.CENTER if c_idx == 3 else (PP_ALIGN.LEFT if c_idx >= 1 else PP_ALIGN.CENTER)
        text_color = fc
        if r_idx > 0 and c_idx == 3:
            text_color = ACCENT
            fb = True
        elif r_idx > 0 and c_idx == 0:
            colors = [MVC_BLUE, MVC_PURPLE, MVC_GREEN]
            text_color = colors[r_idx - 1]
            fb = True
        elif r_idx > 0:
            fb = False
        add_text_box(slide, conf_x[c_idx] + Inches(0.08), ry + Inches(0.08),
                     conf_w[c_idx] - Inches(0.16), conf_rh,
                     val if not (r_idx > 0 and c_idx == 3) else f"✓ {val}",
                     font_size=10, bold=fb, color=text_color, alignment=align)

# Bottom note
add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
             "Separation stricte — aucune logique metier dans les vues, aucun rendu HTML dans les controllers (API JSON uniquement)",
             font_size=10, italic=True, color=GRAY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 7)


# =============================================================================
# SLIDE 8 — Demo Front-office
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 8, "Demo — Parcours citoyen", "Front-office public", "3 min")
# (was slide 6 before MVC + comparaison inserted)

# Large screenshot placeholder
placeholder_image(slide, Inches(0.5), Inches(1.4), Inches(8), Inches(5.5),
                  "DEMO LIVE — Front-office\n\nBasculer vers le navigateur ici\nhttp://localhost:3005")

# Right: checklist of demo steps
steps = [
    "1. Page d'accueil (hero, categories)",
    "2. Listing ressources + filtres",
    "3. Detail ressource + commentaires",
    "4. Inscription nouveau compte",
    "5. Connexion citoyen",
    "6. Ajout favori + commentaire",
    "7. Tableau de progression",
]
add_text_box(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(0.4),
             "Parcours a montrer", font_size=16, bold=True, color=PRIMARY)
add_bullet_list(slide, Inches(8.8), Inches(2.1), Inches(4), Inches(4.5),
                steps, font_size=14, color=DARK_TEXT, spacing=1.5)

# Backup note
add_text_box(slide, Inches(8.8), Inches(6.2), Inches(4), Inches(0.5),
             "Backup : captures d'ecran si Docker ne demarre pas",
             font_size=10, color=GRAY, italic=True)

slide_number_footer(slide, 8)


# =============================================================================
# SLIDE 9 — Demo Back-office
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 9, "Demo — Espace d'administration", "Back-office et moderation", "2 min")

placeholder_image(slide, Inches(0.5), Inches(1.4), Inches(8), Inches(5.5),
                  "DEMO LIVE — Back-office admin\n\nBasculer vers le navigateur ici\nhttp://localhost:3005/administration")

steps_admin = [
    "1. Dashboard stats (barres + donut)",
    "2. Metriques cles (users, actifs...)",
    "3. Moderation : approuver ressource",
    "4. Moderation : gerer commentaires",
    "5. Gestion utilisateurs (CRUD)",
    "6. Gestion categories",
]
add_text_box(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(0.4),
             "Parcours a montrer", font_size=16, bold=True, color=PRIMARY)
add_bullet_list(slide, Inches(8.8), Inches(2.1), Inches(4), Inches(4),
                steps_admin, font_size=14, color=DARK_TEXT, spacing=1.5)

# Login info
info_box = add_shape_rect(slide, Inches(8.8), Inches(5.5), Inches(4), Inches(1.2),
                          RGBColor(0xFD, 0xF5, 0xE6), radius=True)
add_text_box(slide, Inches(9), Inches(5.6), Inches(3.6), Inches(1),
             "Compte demo :\nsuperadmin@ressources.local\nMot de passe : password123",
             font_size=12, color=DARK_TEXT)

slide_number_footer(slide, 9)


# =============================================================================
# SLIDE 10 — Application mobile
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 10, "Application mobile", "React Native / Expo — Cross-platform", "1 min 30")

# Architecture placeholder
placeholder_image(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4),
                  "Schema d'architecture mobile\n\nEcrans React Navigation (stack/tab)\n→ API REST Laravel\n→ AsyncStorage pour le token")

# Mobile screens table
screens = [
    ("Accueil", "Acces rapide ressources", "GET /api/resources"),
    ("Ressources", "Liste + filtres + recherche", "GET /api/resources?..."),
    ("Detail", "Contenu, commentaires, favori", "GET /resources/{id}"),
    ("Connexion", "Authentification", "POST /api/login"),
    ("Profil", "Info utilisateur", "GET /api/user"),
    ("Progression", "Favoris, exploites", "GET /api/progression"),
]

add_text_box(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(0.4),
             "Ecrans principaux", font_size=16, bold=True, color=PRIMARY)

# Mini table
table_top = Inches(2.1)
hdrs = ["Ecran", "Fonctionnalite", "Endpoint API"]
col_xs = [Inches(6.5), Inches(8.3), Inches(10.2)]
col_ws = [Inches(1.8), Inches(1.9), Inches(2.4)]

# Header row
for j, hdr in enumerate(hdrs):
    add_shape_rect(slide, col_xs[j], table_top, col_ws[j], Inches(0.35), PRIMARY)
    add_text_box(slide, col_xs[j] + Inches(0.05), table_top + Inches(0.02),
                 col_ws[j] - Inches(0.1), Inches(0.3),
                 hdr, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

for r_idx, (scr, func, endpoint) in enumerate(screens):
    ry = table_top + Inches(0.35) + r_idx * Inches(0.42)
    bg = LIGHT_BG if r_idx % 2 == 0 else WHITE
    for j, val in enumerate([scr, func, endpoint]):
        add_shape_rect(slide, col_xs[j], ry, col_ws[j], Inches(0.42), bg)
        add_text_box(slide, col_xs[j] + Inches(0.05), ry + Inches(0.05),
                     col_ws[j] - Inches(0.1), Inches(0.32),
                     val, font_size=10, color=DARK_TEXT)

# Key points
add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.8),
             "Un seul backend pour deux clients : l'application mobile consomme exactement la meme API REST que le web — aucune duplication de logique.",
             font_size=14, color=PRIMARY, alignment=PP_ALIGN.CENTER, italic=True)

slide_number_footer(slide, 10)


# =============================================================================
# SLIDE 11 — Securite en profondeur
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 11, "Securite en profondeur", "Couches de protection et headers HTTP", "1 min")

# Security layers diagram (left)
layers = [
    ("Sanitization", "strip_tags, trim, normalisation", ACCENT),
    ("Validation", "FormRequest, regles declaratives", RGBColor(0x42, 0xA5, 0xF5)),
    ("Authentification", "Sanctum, tokens Bearer", PRIMARY),
    ("Controle de roles", "RoleMiddleware + ResourcePolicy", RGBColor(0x7B, 0x1F, 0xA2)),
    ("Chiffrement BDD", "Crypt::encryptString, email_hash SHA-256", RED_SOFT),
]
layer_x = Inches(0.5)
layer_w = Inches(5.8)
layer_h = Inches(0.8)
for i, (name, desc, color) in enumerate(layers):
    ly = Inches(1.5) + i * (layer_h + Inches(0.12))
    add_shape_rect(slide, layer_x, ly, layer_w, layer_h, color, radius=True)
    add_text_box(slide, layer_x + Inches(0.3), ly + Inches(0.08), Inches(2.5), Inches(0.35),
                 name, font_size=15, bold=True, color=WHITE)
    add_text_box(slide, layer_x + Inches(0.3), ly + Inches(0.4), layer_w - Inches(0.6), Inches(0.35),
                 desc, font_size=11, color=RGBColor(0xEE, 0xEE, 0xEE))

# Right: HTTP Headers
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(0.4),
             "Headers de securite HTTP", font_size=16, bold=True, color=PRIMARY)

headers_list = [
    "X-Content-Type-Options: nosniff",
    "X-Frame-Options: DENY",
    "Referrer-Policy: no-referrer",
    "X-XSS-Protection: 1; mode=block",
    "Content-Security-Policy: default-src 'none';\n  frame-ancestors 'none'; base-uri 'self'",
]
add_bullet_list(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(2.5),
                [f"   {h}" for h in headers_list], font_size=12, color=DARK_TEXT)

# Rate limiting box
rl_box = add_shape_rect(slide, Inches(7), Inches(4.6), Inches(5.8), Inches(1),
                        RGBColor(0xFD, 0xEC, 0xEC), radius=True)
add_text_box(slide, Inches(7.3), Inches(4.7), Inches(5.2), Inches(0.8),
             "Rate Limiting : 5 tentatives/min sur /login et /register\nHTTP 429 Too Many Requests apres depassement\nTeste dans AuthEndpointsTest",
             font_size=12, color=RED_SOFT)

# RGPD box
rgpd_box = add_shape_rect(slide, Inches(7), Inches(5.8), Inches(5.8), Inches(0.9),
                          RGBColor(0xE8, 0xF5, 0xE9), radius=True)
add_text_box(slide, Inches(7.3), Inches(5.9), Inches(5.2), Inches(0.7),
             "RGPD : anonymisation logique via DELETE /api/user\nNom → 'Deleted User' | Email → deleted-user-{uuid}@example.invalid",
             font_size=12, color=ACCENT)

slide_number_footer(slide, 11)


# =============================================================================
# SLIDE 12 — Organisation
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 12, "Notre organisation", "Sprints, Git Flow et conventions", "1 min")

# Sprint timeline
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5), Inches(0.4),
             "Timeline des sprints", font_size=16, bold=True, color=PRIMARY)

sprints = [
    ("S1", "Fondations", "Setup Laravel,\nNext.js, Expo, Auth", ACCENT),
    ("S2", "Coeur", "CRUD ressources,\ncommentaires, favoris", RGBColor(0x42, 0xA5, 0xF5)),
    ("S3", "Avance", "Admin, stats,\nRGAA, progression", RGBColor(0x7B, 0x1F, 0xA2)),
    ("S4", "Polish", "Tests, docs,\ncorrections, slides", ACCENT_GOLD),
]

sprint_w = Inches(2.8)
sprint_h = Inches(1.6)
sprint_gap = Inches(0.25)
sprint_y = Inches(2.1)

for i, (code, name, desc, color) in enumerate(sprints):
    sx = Inches(0.5) + i * (sprint_w + sprint_gap)
    card = add_shape_rect(slide, sx, sprint_y, sprint_w, sprint_h, color, radius=True)
    add_text_box(slide, sx, sprint_y + Inches(0.15), sprint_w, Inches(0.35),
                 f"{code} — {name}", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, sx + Inches(0.2), sprint_y + Inches(0.55), sprint_w - Inches(0.4), Inches(0.9),
                 desc, font_size=12, color=RGBColor(0xEE, 0xEE, 0xEE), alignment=PP_ALIGN.CENTER)
    # Arrow between sprints
    if i < 3:
        arrow_x = sx + sprint_w + Inches(0.02)
        add_text_box(slide, arrow_x, sprint_y + Inches(0.5), Inches(0.25), Inches(0.4),
                     "→", font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

# Git Flow
add_text_box(slide, Inches(0.6), Inches(4.2), Inches(6), Inches(0.4),
             "Git Flow", font_size=16, bold=True, color=PRIMARY)

git_lines = [
    "main         Production stable, merge en fin de sprint",
    "develop     Integration, toutes les features mergees via PR",
    "feature/*   Branches par fonctionnalite, squash and merge",
]
add_bullet_list(slide, Inches(0.8), Inches(4.7), Inches(6), Inches(1.5),
                [f"   {l}" for l in git_lines], font_size=13, color=DARK_TEXT)

# Conventions
add_text_box(slide, Inches(7), Inches(4.2), Inches(5.8), Inches(0.4),
             "Conventional Commits", font_size=16, bold=True, color=PRIMARY)

commits = [
    "feat:      nouvelle fonctionnalite",
    "fix:        correction de bug",
    "docs:     documentation",
    "test:      ajout de tests",
    "refactor: refactorisation",
]
add_bullet_list(slide, Inches(7.2), Inches(4.7), Inches(5.5), Inches(1.8),
                [f"   {c}" for c in commits], font_size=13, color=DARK_TEXT)

# Rule box
rule = add_shape_rect(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.7),
                      RGBColor(0xFD, 0xF5, 0xE6), radius=True)
add_text_box(slide, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
             "Regle d'or : jamais de push direct sur main ou develop — tout passe par Pull Request avec review obligatoire",
             font_size=13, bold=True, color=RGBColor(0xCC, 0x77, 0x00), alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 12)


# =============================================================================
# SLIDE 13 — Bilan
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, WHITE)
section_header(slide, 13, "Bilan", "Ce qui a bien fonctionne", "1 min")

strengths = [
    ("Architecture MVC stricte", "Laravel : maintenabilite, testabilite,\nseparation claire des responsabilites", "🏗️"),
    ("Securite renforcee", "Chiffrement donnees, anonymisation RGPD,\nheaders HTTP, rate limiting", "🔒"),
    ("Accessibilite RGAA", "Skip links, ARIA, focus visible, contraste AAA\nintegres des la conception", "♿"),
    ("27 tests backend + 117 tests frontend", "Backend : PHPUnit (Auth, CRUD, admin, moderation)\nFrontend : Vitest + Testing Library (API, composants, contextes)", "✅"),
    ("Docker reproductible", "4 services containerises, health checks,\nenvironnement identique pour tous", "🐳"),
    ("Workflow Git structure", "PR obligatoires, review, conventional commits,\nsquash and merge", "🔀"),
]

card_w = Inches(3.85)
card_h = Inches(1.55)
gap_x = Inches(0.4)
gap_y = Inches(0.3)

for i, (title, desc, icon) in enumerate(strengths):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * (card_w + gap_x)
    y = Inches(1.5) + row * (card_h + gap_y)
    add_card(slide, x, y, card_w, card_h, title, desc, icon=icon, bg_color=LIGHT_BG)

slide_number_footer(slide, 13)


# =============================================================================
# SLIDE 14 — Perspectives
# =============================================================================

slide = prs.slides.add_slide(prs.slide_layouts[6])
gradient_bg(slide, PRIMARY_DARK, PRIMARY)

# Title
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             "Et si on continuait ?", font_size=32, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)
add_text_box(slide, Inches(0.6), Inches(1), Inches(12), Inches(0.4),
             "Perspectives d'evolution", font_size=16, color=RGBColor(0xAA, 0xBB, 0xDD),
             alignment=PP_ALIGN.CENTER, italic=True)

# Three columns: court / moyen / long terme
phases = [
    ("Court terme", [
        "Tests E2E\n(Cypress / Playwright)",
        "Enrichissement app mobile",
        "Couverture de code\n(Istanbul / c8)",
    ], ACCENT),
    ("Moyen terme", [
        "Messagerie entre\nparticipants",
        "Export statistiques\n(CSV, PDF)",
        "Pipeline CI/CD\n(GitHub Actions)",
    ], RGBColor(0x42, 0xA5, 0xF5)),
    ("Long terme", [
        "Support multilingue",
        "Notifications push\nmobile",
        "Mode hors-ligne\npartiel",
    ], ACCENT_GOLD),
]

col_w = Inches(3.8)
col_gap = Inches(0.35)
col_y = Inches(1.7)

for i, (phase, items, color) in enumerate(phases):
    cx = Inches(0.6) + i * (col_w + col_gap)
    # Phase header
    hdr = add_shape_rect(slide, cx, col_y, col_w, Inches(0.55), color, radius=True)
    add_text_box(slide, cx, col_y + Inches(0.08), col_w, Inches(0.4),
                 phase, font_size=16, bold=True, color=WHITE if i != 2 else DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)
    # Items
    for j, item in enumerate(items):
        iy = col_y + Inches(0.7) + j * Inches(1.35)
        item_card = add_shape_rect(slide, cx + Inches(0.1), iy, col_w - Inches(0.2), Inches(1.15),
                                   RGBColor(0x15, 0x2E, 0x55), radius=True)
        item_card.line.color.rgb = RGBColor(0x3A, 0x5A, 0x8A)
        item_card.line.width = Pt(1)
        add_text_box(slide, cx + Inches(0.3), iy + Inches(0.15), col_w - Inches(0.6), Inches(0.85),
                     item, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

# Bottom
add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             "Merci de votre attention — Place aux questions",
             font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 14)


# =============================================================================
# SAVE
# =============================================================================

prs.save(OUTPUT)
print(f"PowerPoint genere : {OUTPUT}")
print(f"Nombre de slides : {len(prs.slides)}")
