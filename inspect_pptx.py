"""Inspecte les slides 7, 10, 13 pour voir la structure des paragraphes restants."""
from pptx import Presentation
from pathlib import Path

prs = Presentation(str(Path(__file__).parent / "Soutenance_RE_Sources_NEW.pptx"))

TARGETS = {
    7:  ["Category", "Moderation", "Favorite", "Progression", "Resource"],
    10: ["AsyncStorage", "secure-store", "Schema d'architecture", "Expo Router", "React Navigation"],
    13: ["Vitest", "PHPUnit", "Testing Library", "Frontend"],
}

def walk(shapes):
    for s in shapes:
        if s.has_text_frame:
            yield s.text_frame
        if s.shape_type == 6:
            yield from walk(s.shapes)

for idx, slide in enumerate(prs.slides, 1):
    if idx not in TARGETS:
        continue
    print(f"\n=== SLIDE {idx} ===")
    for tf_i, tf in enumerate(walk(slide.shapes)):
        for p_i, p in enumerate(tf.paragraphs):
            runs = [r.text for r in p.runs]
            full = "".join(runs)
            for kw in TARGETS[idx]:
                if kw in full:
                    safe = full.encode('ascii', 'replace').decode('ascii')
                    print(f"  tf#{tf_i} p#{p_i} | runs={len(runs)} | {safe!r}")
                    break
