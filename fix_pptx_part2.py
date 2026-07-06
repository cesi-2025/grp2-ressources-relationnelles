"""
Deuxieme passe : corrige les textes qui etaient decoupes sur plusieurs runs
et que fix_pptx.py n'a pas pu atteindre.
"""
from pptx import Presentation
from pathlib import Path

SRC = Path(__file__).parent / "Soutenance_RE_Sources_NEW.pptx"
prs = Presentation(str(SRC))


def walk(shapes):
    for s in shapes:
        if s.has_text_frame:
            yield s.text_frame
        if s.shape_type == 6:
            yield from walk(s.shapes)


def set_paragraph_text(para, new_text: str):
    """Ecrase le texte du paragraphe en gardant le format du 1er run."""
    if not para.runs:
        return False
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ""
    return True


count = 0

# ---------------- SLIDE 7 ----------------
slide7 = prs.slides[6]
for tf in walk(slide7.shapes):
    for para in tf.paragraphs:
        full = "".join(r.text for r in para.runs)
        if "+ Category, Progression," in full and "RelationType" in full:
            set_paragraph_text(para,
                "+ Category, Progression, RelationType, "
                "ResourceType, ActivitySession, ResourceView, SearchLog")
            print(f"  slide 7 : liste Models etendue")
            count += 1
        elif "+ Moderation, Favorite," in full and "Progression" in full:
            set_paragraph_text(para,
                "+ Moderation, Favorite, Progression, "
                "SuperAdmin, Activity, Category")
            print(f"  slide 7 : liste Controllers etendue")
            count += 1

# ---------------- SLIDE 10 ----------------
slide10 = prs.slides[9]
for tf in walk(slide10.shapes):
    for para in tf.paragraphs:
        runs = list(para.runs)
        for run in runs:
            if "AsyncStorage pour le token" in run.text:
                run.text = run.text.replace(
                    "AsyncStorage pour le token",
                    "expo-secure-store (Keychain / Keystore)")
                print(f"  slide 10 : AsyncStorage -> expo-secure-store")
                count += 1
            if "Ecrans React Navigation (stack/tab)" in run.text:
                run.text = run.text.replace(
                    "Ecrans React Navigation (stack/tab)",
                    "Expo Router (file-based) + Stack natif")
                print(f"  slide 10 : navigation corrigee")
                count += 1

# ---------------- SLIDE 13 ----------------
slide13 = prs.slides[12]
for tf in walk(slide13.shapes):
    for para in tf.paragraphs:
        full = "".join(r.text for r in para.runs)
        if "Vitest" in full and "Frontend" in full:
            set_paragraph_text(para,
                "Backend : PHPUnit (Auth, CRUD, admin, moderation, interactions)\n"
                "Mobile : Jest + @testing-library/react-native (services)")
            print(f"  slide 13 : description tests corrigee")
            count += 1

prs.save(str(SRC))
print(f"\nTotal 2e passe : {count} corrections")
