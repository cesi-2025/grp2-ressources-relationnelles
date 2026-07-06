"""Verifie qu'il ne reste plus aucune des chaines erronees."""
from pptx import Presentation
from pathlib import Path

prs = Presentation(str(Path(__file__).parent / "Soutenance_RE_Sources_NEW.pptx"))

BAD = [
    "Laravel 11 (retenu)",
    "8 modeles Eloquent",
    "9 controllers",
    "/administration",
    "AsyncStorage pour le token",
    "React Navigation (stack/tab)",
    "27 tests backend",
    "117 tests frontend",
    "Vitest + Testing Library",
]

def walk(shapes):
    for s in shapes:
        if s.has_text_frame: yield s.text_frame
        if s.shape_type == 6: yield from walk(s.shapes)

found = []
for idx, slide in enumerate(prs.slides, 1):
    for tf in walk(slide.shapes):
        for para in tf.paragraphs:
            full = "".join(r.text for r in para.runs)
            for bad in BAD:
                if bad in full:
                    found.append((idx, bad, full[:80]))

if not found:
    print("OK : aucune chaine erronee restante.")
else:
    print("RESIDUS detectes :")
    for idx, bad, ctx in found:
        print(f"  slide {idx} | '{bad}' | {ctx!r}")

# Check qu'on a bien le nouveau contenu
GOOD = [
    ("slide 5 Laravel 12",   "Laravel 12 (retenu)"),
    ("slide 7 13 modeles",   "13 modeles Eloquent"),
    ("slide 7 10 controllers", "10 controllers"),
    ("slide 9 /admin",       "localhost:3005/admin"),
    ("slide 10 secure-store", "expo-secure-store"),
    ("slide 13 30 tests",    "30 tests backend + 14 tests mobile"),
]
print("\nVerifications positives :")
for label, text in GOOD:
    ok = False
    for slide in prs.slides:
        for tf in walk(slide.shapes):
            for para in tf.paragraphs:
                if text in "".join(r.text for r in para.runs):
                    ok = True; break
            if ok: break
        if ok: break
    print(f"  [{'OK' if ok else 'KO'}] {label} ('{text}')")
