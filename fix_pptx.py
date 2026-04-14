"""
Corrige Soutenance_RE_Sources_NEW.pptx sans regenerer : modifie uniquement
les chaines incorrectes, preserve toutes les images ajoutees manuellement.

Corrections basees sur audit vs code reel au 2026-04-13.
"""

from pptx import Presentation
from pathlib import Path
import shutil

ROOT = Path(__file__).parent
SRC = ROOT / "Soutenance_RE_Sources_NEW.pptx"
BACKUP = ROOT / "Soutenance_RE_Sources_NEW.backup.pptx"
OUT = SRC  # on ecrase apres backup

# Backup
if not BACKUP.exists():
    shutil.copy2(SRC, BACKUP)
    print(f"[backup] {BACKUP.name}")

# Chaque entree: (ancien_texte, nouveau_texte)
# Ordre important : les remplacements longs doivent passer avant les courts
# pour eviter les conflits de substring.
REPLACEMENTS = [
    # -------------------------------------------------------------------
    # SLIDE 5 — Laravel 11 -> Laravel 12
    # -------------------------------------------------------------------
    ("Laravel 11 (retenu)", "Laravel 12 (retenu)"),

    # -------------------------------------------------------------------
    # SLIDE 7 — MVC : modeles, controllers, chemins
    # -------------------------------------------------------------------
    # Tableau conformite MVC (lignes)
    ("8 modeles Eloquent (User, Resource, Comment...)",
     "13 modeles Eloquent (User, Resource, Comment, Activity...)"),
    ("9 controllers + FormRequest + Policies",
     "10 controllers + FormRequest + Policies"),

    # Liste fichiers Model (colonne MODEL)
    ("+ Category, Progression,\n  RelationType, ResourceType",
     "+ Category, Progression, RelationType,\n  ResourceType, ActivitySession,\n  ResourceView, SearchLog"),

    # Liste fichiers Controller (colonne CONTROLLER)
    ("+ Moderation, Favorite,\n  Progression, SuperAdmin",
     "+ Moderation, Favorite, Progression,\n  SuperAdmin, Activity, Category"),

    # Liste fichiers View (colonne VIEW) — administration/ -> admin/
    ("(main)/auth/connexion/", "(main)/auth/connexion/"),  # no-op garde-fou
    ("administration/", "admin/"),

    # -------------------------------------------------------------------
    # SLIDE 9 — URL back-office
    # -------------------------------------------------------------------
    ("http://localhost:3005/administration",
     "http://localhost:3005/admin"),

    # -------------------------------------------------------------------
    # SLIDE 10 — Application mobile
    # -------------------------------------------------------------------
    # Schema placeholder : AsyncStorage -> expo-secure-store, navigation
    ("Schema d'architecture mobile\n\nEcrans React Navigation (stack/tab)\n"
     "→ API REST Laravel\n→ AsyncStorage pour le token",
     "Schema d'architecture mobile\n\nExpo Router (file-based) + Stack natif\n"
     "→ API REST Laravel (JWT Sanctum)\n→ expo-secure-store (Keychain/Keystore)"),

    # Ecrans mobiles : aligner sur le code reel
    # Slide affiche : Accueil / Ressources / Detail / Connexion / Profil / Progression
    # Reel : index, login, register, dashboard, profile, favorites,
    #        resource/[id], resource/create, resource/edit
    ("Acces rapide ressources", "Liste ressources + filtres"),
    ("Liste + filtres + recherche", "Creation ressource"),
    ("GET /api/resources?...", "POST /api/resources"),
    ("Contenu, commentaires, favori", "Detail + commentaires + favori"),
    ("Info utilisateur", "Profil + deconnexion"),
    ("Favoris, exploites", "Liste des favoris"),
    ("GET /api/progression", "GET /api/resources/{id}/favorite"),
    # Renommage premiere colonne
    # (Accueil reste ; Ressources -> Creation ; Progression -> Favoris)

    # -------------------------------------------------------------------
    # SLIDE 13 — Bilan : nombre de tests
    # -------------------------------------------------------------------
    ("27 tests backend + 117 tests frontend",
     "30 tests backend + 14 tests mobile"),
    ("Backend : PHPUnit (Auth, CRUD, admin, moderation)\n"
     "Frontend : Vitest + Testing Library (API, composants, contextes)",
     "Backend : PHPUnit (Auth, CRUD, admin, moderation, interactions)\n"
     "Mobile : Jest + @testing-library/react-native (couche services)"),
]

# Corrections supplementaires pour slide 10 : renommer la colonne "Ecran"
# (il faut des remplacements plus cibles car "Ressources" et "Progression"
# sont ambigus). On les traite apres.
SLIDE10_ROW_FIXES = [
    # (ancien_label, nouveau_label)  — ligne par ligne dans l'ordre d'apparition
    ("Ressources", "Creation"),
    ("Progression", "Favoris"),
]


def replace_in_runs(text_frame, old: str, new: str) -> int:
    """Remplace old->new dans un text_frame.
    Gere le cas ou old est decoupe sur plusieurs runs en concatenant
    puis en reinjectant dans le premier run non vide."""
    count = 0
    for para in text_frame.paragraphs:
        # 1) tentative simple run-par-run
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                count += 1
        # 2) tentative paragraphe entier si non trouve
        full = "".join(r.text for r in para.runs)
        if old in full and old not in "".join(r.text for r in para.runs if old in r.text):
            # old traverse plusieurs runs
            new_full = full.replace(old, new)
            if para.runs:
                para.runs[0].text = new_full
                for r in para.runs[1:]:
                    r.text = ""
                count += 1
    return count


def walk_shapes(shapes):
    for shape in shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.shape_type == 6:  # GROUP
            yield from walk_shapes(shape.shapes)


def apply_replacements(prs):
    total = 0
    for idx, slide in enumerate(prs.slides, start=1):
        for tf in walk_shapes(slide.shapes):
            for old, new in REPLACEMENTS:
                if old == new:
                    continue
                n = replace_in_runs(tf, old, new)
                if n:
                    total += n
                    snippet = old[:50].replace("\n", " / ")
                    print(f"  slide {idx:>2} : '{snippet}...' x{n}")
    return total


def fix_slide10_screens(prs):
    """Slide 10 : tableau des ecrans mobiles. Renomme 'Ressources' -> 'Creation'
    et 'Progression' -> 'Favoris' uniquement dans ce slide."""
    if len(prs.slides) < 10:
        return 0
    slide = prs.slides[9]  # slide 10 (index 9)
    changes = 0
    for tf in walk_shapes(slide.shapes):
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text.strip() == "Ressources":
                    run.text = "Creation"
                    changes += 1
                elif run.text.strip() == "Progression":
                    run.text = "Favoris"
                    changes += 1
    return changes


def main():
    prs = Presentation(str(SRC))
    print(f"Ouverture : {SRC.name}  ({len(prs.slides)} slides)")

    print("\n[corrections globales]")
    n_global = apply_replacements(prs)

    print("\n[corrections slide 10 : tableau ecrans]")
    n_s10 = fix_slide10_screens(prs)
    print(f"  {n_s10} cellules renommees (Ressources->Creation, Progression->Favoris)")

    prs.save(str(OUT))
    print(f"\nTotal remplacements : {n_global + n_s10}")
    print(f"Sauvegarde : {OUT.name}")
    print(f"Backup original : {BACKUP.name}")


if __name__ == "__main__":
    main()
